from pptx import Presentation

# Step 1: Extract Text from PowerPoint Slides
def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    slides_text = []

    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text
                slide_text.append(text)
        slides_text.append(slide_text)
    
    return slides_text

# Example usage:
file_path = r"C:\Users\vipas\Downloads\Minimalist sales pitch.pptx"
slides_text = extract_text_from_ppt(file_path)

# Display extracted text
for i, slide_text in enumerate(slides_text):
    print(f"Slide {i+1}:")
    for text in slide_text:
        print(text)
    print()

# Step 2: Prepare New Text for Each Text Box
new_text = [
    ["New text for slide 1, text box 1", "New text for slide 1, text box 2"],
    ["New text for slide 2, text box 1", "New text for slide 2, text box 2"],
    # Add more slides as needed
]

# Step 3: Update the PowerPoint File with Modified Text
def update_ppt_with_modified_text(original_file_path, modified_text, output_file_path):
    prs = Presentation(original_file_path)

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                if i < len(modified_text) and j < len(modified_text[i]):
                    text_frame.text = modified_text[i][j]

    prs.save(output_file_path)

# Example usage:
output_file_path = r"C:\Users\vipas\Downloads\modified_presentation1.pptx"
update_ppt_with_modified_text(file_path, new_text, output_file_path)
